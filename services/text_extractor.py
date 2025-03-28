import json
import PyPDF2
from pptx import Presentation
from services.ocr_service import OCRService
from services.ai_service import AIService
from utils.nlp_utils import parse_text_with_nlp
from config import Config

class TextExtractor:
    
    @staticmethod
    def extract_from_pdf(file_path):
        """Extract text from a PDF file using PyPDF2, with OCR fallback."""
        try:
            text = ""  # Initialize text variable
            metadata = {}  # Initialize metadata variable

            # Read PDF file before closing it
            with open(file_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
                metadata = dict(reader.metadata) if reader.metadata else {}

            # If no text found, use OCR fallback
            if not text.strip():
                text = AIService.extract_pitch_deck_info(file_path)
            categorized_sections = parse_text_with_nlp(text, Config.SECTION_KEYWORDS)

            return categorized_sections

        except Exception as e:
            return f"Error extracting PDF: {str(e)}"

        
        

    @staticmethod
    def extract_from_pptx(file_path):
        """Extract text from a PowerPoint file."""
        try:
            prs = Presentation(file_path)
            sections = {}
            for i, slide in enumerate(prs.slides):
                title = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
                content = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                sections[title] = content
            return sections
        except Exception as e:
            return f"Error extracting PPTX: {str(e)}"



